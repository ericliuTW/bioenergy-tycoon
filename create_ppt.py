"""
生質能校園大亨 — 教師活動後總結簡報 產生器
Forest & Moss palette
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = r'D:\caludecode\bioenergy-tycoon\bioenergy-teacher-summary.pptx'

# ── Color palette ──
C_DARK    = RGBColor(0x14, 0x1E, 0x14)
C_PRIMARY = RGBColor(0x2C, 0x5F, 0x2D)
C_SEC     = RGBColor(0x97, 0xBC, 0x62)
C_ACCENT  = RGBColor(0xF5, 0xF5, 0xF0)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x1E, 0x1E, 0x1E)
C_GRAY    = RGBColor(0x6B, 0x6B, 0x6B)
C_LTGRAY  = RGBColor(0xE8, 0xE8, 0xE0)
C_ORANGE  = RGBColor(0xE8, 0x8D, 0x2A)
C_TEAL    = RGBColor(0x1A, 0x93, 0x6F)

FONT_TITLE = '微軟正黑體'
FONT_BODY  = '微軟正黑體'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ── Helper functions ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, x, y, w, h, text, font_size=16, color=C_BLACK, bold=False,
                 align=PP_ALIGN.LEFT, font_name=FONT_BODY, line_space=1.3):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(0)
    if line_space != 1.0:
        p.line_spacing = Pt(font_size * line_space)
    tf.auto_size = None
    return txBox

def add_multiline_box(slide, x, y, w, h, lines, font_size=14, color=C_BLACK, bold=False,
                      align=PP_ALIGN.LEFT, font_name=FONT_BODY, line_space=1.5):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(6)
        if line_space != 1.0:
            p.line_spacing = Pt(font_size * line_space)
    return txBox

def card(slide, x, y, w, h, title, body_lines, icon='', title_color=C_PRIMARY):
    add_rect(slide, x, y, w, h, C_WHITE)
    add_rect(slide, x, y, Inches(0.07), h, title_color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.15), w - Inches(0.4), Inches(0.5),
                 f'{icon} {title}' if icon else title, font_size=18, color=title_color, bold=True)
    add_multiline_box(slide, x + Inches(0.25), y + Inches(0.65), w - Inches(0.4), h - Inches(0.75),
                      body_lines, font_size=14, color=C_BLACK, line_space=1.5)

def stat_block(slide, x, y, number, label, color=C_SEC):
    add_text_box(slide, x, y, Inches(2.5), Inches(0.8), str(number),
                 font_size=44, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.75), Inches(2.5), Inches(0.5), label,
                 font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

def section_header(slide, num, title):
    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(1), Inches(0.6),
                 num, font_size=38, color=C_SEC, bold=True)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(10), Inches(0.6),
                 title, font_size=30, color=C_PRIMARY, bold=True)
    add_rect(slide, Inches(0.6), Inches(1.1), Inches(5), Inches(0.05), C_SEC)


# ════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), C_SEC)

add_text_box(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(1.2),
             '🌱 生質能校園大亨：微電網保衛戰',
             font_size=46, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, font_name=FONT_TITLE)

add_text_box(slide, Inches(1), Inches(2.9), Inches(11.3), Inches(0.7),
             'Bioenergy Campus Tycoon: Microgrid Defense',
             font_size=22, color=C_SEC, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.05), C_SEC)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.7),
             '教師活動後總結簡報',
             font_size=28, color=C_LTGRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6),
             '適用對象：高中學生  |  活動時間：15 分鐘遊戲 + 25 分鐘討論',
             font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), SH - Inches(0.08), SW, Inches(0.08), C_SEC)


# ════════════════════════════════════════════
# SLIDE 2: 活動回顧
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
section_header(slide, '01', '活動回顧 — 你們剛剛做了什麼？')

# 4 cards in 2x2
card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(1.6),
     '收集與前處理', [
         '校園廚餘、枯枝、畜牧糞水 → 破碎、去砂、混合',
         '調整 C/N 比到最佳範圍（20-30）'
     ], icon='🗑️')

card(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(1.6),
     '厭氧消化與氣化', [
         '消化槽：有機物 → 沼氣（CH4 + CO2）',
         '氣化爐：乾燥生物質 → 合成氣'
     ], icon='🔬')

card(slide, Inches(0.6), Inches(3.4), Inches(5.8), Inches(1.6),
     '發電與供電', [
         '沼氣/合成氣 → 發電機 → 電力',
         '拉電線到校園建築，建立微電網'
     ], icon='⚡')

card(slide, Inches(6.9), Inches(3.4), Inches(5.8), Inches(1.6),
     '外部成本管理', [
         '臭味、噪音、設備距離、砍樹抗議',
         '洗滌塔可降低排放，但需要資金'
     ], icon='🏘️')

# Bottom highlight
add_rect(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.1), C_PRIMARY)
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(0.7),
             '💡 討論：你的系統最終成績如何？你覺得最大的挑戰是什麼？',
             font_size=20, color=C_WHITE, bold=True)


# ════════════════════════════════════════════
# SLIDE 3: 核心知識 — 厭氧消化（拆成兩頁）
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
section_header(slide, '02', '核心知識 — 厭氧消化原理')

# Left: Process flow
add_rect(slide, Inches(0.6), Inches(1.5), Inches(6.2), Inches(5.5), C_WHITE)
add_rect(slide, Inches(0.6), Inches(1.5), Inches(0.07), Inches(5.5), C_PRIMARY)

add_text_box(slide, Inches(1.0), Inches(1.6), Inches(5.5), Inches(0.5),
             '🔄 四階段反應流程', font_size=20, color=C_PRIMARY, bold=True)

steps = [
    ('水解 Hydrolysis', '大分子有機物 → 單糖、胺基酸、脂肪酸'),
    ('酸化 Acidogenesis', '單糖等 → 揮發性脂肪酸（VFA）+ CO2 + H2'),
    ('乙酸化 Acetogenesis', 'VFA → 乙酸 + CO2 + H2'),
    ('甲烷化 Methanogenesis', '乙酸 + CO2 + H2 → CH4 + CO2（沼氣）')
]
y_pos = Inches(2.3)
for i, (title, desc) in enumerate(steps):
    # Step number circle
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y_pos, Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = C_PRIMARY
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i+1)
    p.font.size = Pt(16)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.85), y_pos - Inches(0.03), Inches(4.7), Inches(0.4),
                 title, font_size=16, color=C_BLACK, bold=True)
    add_text_box(slide, Inches(1.85), y_pos + Inches(0.38), Inches(4.7), Inches(0.4),
                 desc, font_size=14, color=C_GRAY)

    if i < 3:
        add_rect(slide, Inches(1.38), y_pos + Inches(0.5), Inches(0.04), Inches(0.45), C_SEC)

    y_pos += Inches(1.2)

# Right top: Key parameters
add_rect(slide, Inches(7.1), Inches(1.5), Inches(5.6), Inches(2.6), C_WHITE)
add_rect(slide, Inches(7.1), Inches(1.5), Inches(0.07), Inches(2.6), C_ORANGE)
add_text_box(slide, Inches(7.4), Inches(1.6), Inches(5.1), Inches(0.5),
             '📊 關鍵參數', font_size=20, color=C_ORANGE, bold=True)
params = [
    'C/N 比：最佳 20-30',
    'pH 值：維持 6.8-7.4',
    '有機負荷率（OLR）：進料速度控制',
    '溫度：中溫消化 35-40°C'
]
add_multiline_box(slide, Inches(7.4), Inches(2.2), Inches(5.1), Inches(1.8),
                  params, font_size=14, color=C_BLACK, line_space=1.6)

# Right bottom: Game connection
add_rect(slide, Inches(7.1), Inches(4.4), Inches(5.6), Inches(2.6), C_WHITE)
add_rect(slide, Inches(7.1), Inches(4.4), Inches(0.07), Inches(2.6), C_TEAL)
add_text_box(slide, Inches(7.4), Inches(4.5), Inches(5.1), Inches(0.5),
             '🎮 遊戲中的對應', font_size=20, color=C_TEAL, bold=True)
game_items = [
    '混合槽決定 C/N 比 → 影響產氣量',
    'pH 過低 = 酸化事故 → 需花錢修復',
    '進料太快 → 穩定度下降、系統崩潰',
    '脫硫器 → 保護發電機、降低排放'
]
add_multiline_box(slide, Inches(7.4), Inches(5.1), Inches(5.1), Inches(1.8),
                  game_items, font_size=14, color=C_BLACK, line_space=1.6)


# ════════════════════════════════════════════
# SLIDE 4: 外部成本與社會影響
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
section_header(slide, '03', '外部成本與社會影響')

# Externality definition
add_rect(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.2), C_WHITE)
add_rect(slide, Inches(0.6), Inches(1.4), Inches(0.07), Inches(1.2), C_ORANGE)
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11.4), Inches(0.45),
             '🔑 什麼是外部成本（Externality）？', font_size=20, color=C_ORANGE, bold=True)
add_text_box(slide, Inches(1.0), Inches(1.95), Inches(11.4), Inches(0.55),
             '生產或消費行為對第三方造成的成本，但未反映在市場價格中。',
             font_size=16, color=C_BLACK)

# 4 types - 2x2 grid with more space
types = [
    ('💨 臭味排放', '消化槽與處理設備的氣味擴散', '抗議值上升，嚴重時被迫停工'),
    ('🔊 噪音汙染', '發電機、破碎機運轉噪音影響教學', '設備數量越多，噪音成本越高'),
    ('📏 距離過近', '設備離教室太近造成安全問題', '設備與校園建築距離影響抗議值'),
    ('🌳 砍樹抗議', '移除校園綠地引發師生反彈', '在樹木綠地上建設額外增加抗議')
]
for i, (title, real, game) in enumerate(types):
    col = i % 2
    row = i // 2
    cx = Inches(0.6 + col * 6.3)
    cy = Inches(2.9 + row * 1.65)
    add_rect(slide, cx, cy, Inches(5.8), Inches(1.45), C_WHITE)
    add_rect(slide, cx, cy, Inches(0.07), Inches(1.45), [C_ORANGE, C_TEAL, C_PRIMARY, C_SEC][i])
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.1), Inches(5.3), Inches(0.4),
                 title, font_size=18, color=C_BLACK, bold=True)
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.5), Inches(5.3), Inches(0.4),
                 f'現實：{real}', font_size=14, color=C_GRAY)
    add_text_box(slide, cx + Inches(0.25), cy + Inches(0.9), Inches(5.3), Inches(0.4),
                 f'遊戲：{game}', font_size=14, color=C_TEAL)

# Discussion
add_rect(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.85), C_PRIMARY)
add_text_box(slide, Inches(1.2), Inches(6.45), Inches(11), Inches(0.6),
             '💡 討論：鄰避效應（NIMBY）— 你支持在「你家附近」蓋生質能設施嗎？',
             font_size=20, color=C_WHITE, bold=True)


# ════════════════════════════════════════════
# SLIDE 5: 微電網與能源韌性
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
section_header(slide, '04', '微電網與能源韌性')

# Left: Microgrid
add_rect(slide, Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.8), C_WHITE)
add_rect(slide, Inches(0.6), Inches(1.4), Inches(0.07), Inches(2.8), C_PRIMARY)
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5.2), Inches(0.5),
             '🔋 什麼是微電網？', font_size=20, color=C_PRIMARY, bold=True)
add_multiline_box(slide, Inches(1.0), Inches(2.1), Inches(5.2), Inches(2.0), [
    '小規模、可獨立運作的電力系統',
    '平時與主電網並聯，災害時可「孤島運轉」',
    '整合分散式發電（太陽能、生質能、儲能）',
    '台灣案例：綠島微電網、澎湖智慧電網'
], font_size=15, color=C_BLACK, line_space=1.6)

# Right: Resilience
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.8), C_WHITE)
add_rect(slide, Inches(6.8), Inches(1.4), Inches(0.07), Inches(2.8), C_ORANGE)
add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5.4), Inches(0.5),
             '🛡️ 能源韌性（Energy Resilience）', font_size=20, color=C_ORANGE, bold=True)
add_multiline_box(slide, Inches(7.1), Inches(2.1), Inches(5.4), Inches(2.0), [
    '韌性 = 系統承受衝擊並恢復的能力',
    '極端氣候（颱風、地震）→ 大規模停電風險',
    '2021 台灣 513/517 大停電影響數百萬戶',
    '關鍵設施（醫院、學校）需要備援電力'
], font_size=15, color=C_BLACK, line_space=1.6)

# Game mapping
add_rect(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(1.7), C_WHITE)
add_rect(slide, Inches(0.6), Inches(4.5), Inches(0.07), Inches(1.7), C_TEAL)
add_text_box(slide, Inches(1.0), Inches(4.6), Inches(11.4), Inches(0.5),
             '🎮 遊戲中的設計邏輯', font_size=20, color=C_TEAL, bold=True)
add_multiline_box(slide, Inches(1.0), Inches(5.15), Inches(11.4), Inches(1.0), [
    '平時接電微量加分 → 反映日常備援的「保險」價值',
    '停電時接電大量加分 → 災害時微電網價值才真正顯現',
    '「必接任務建築」→ 模擬現實中優先保障的關鍵設施'
], font_size=15, color=C_BLACK, line_space=1.5)

# Discussion
add_rect(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.85), C_PRIMARY)
add_text_box(slide, Inches(1.2), Inches(6.55), Inches(11), Inches(0.6),
             '💡 討論：如果你是校長，停電時最優先供電的三棟建築是什麼？',
             font_size=20, color=C_WHITE, bold=True)


# ════════════════════════════════════════════
# SLIDE 6: 計分邏輯
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
section_header(slide, '05', '計分邏輯 — 永續發展的多面向評估')

# Scoring table
categories = [
    ('能源產出',   '25%', '總發電量與供電穩定性'),
    ('系統韌性',   '20%', '停電時的供電能力'),
    ('經濟效率',   '10%', '資金運用與維護成本控管'),
    ('系統穩定',   '10%', '消化槽穩定度與 pH 控制'),
    ('時間獎勵',   '10%', '提前完成設計的效率獎勵'),
    ('外部成本',   '5%',  '抗議值控制'),
    ('物流效率',   '5%',  '管線效率與路徑規劃'),
    ('電線覆蓋',   '5%',  '校園建築接電覆蓋率'),
    ('資源利用',   '5%',  '原料使用率與多樣性'),
    ('C/N 優化',   '5%',  'C/N 比控制在最佳範圍'),
]

table_x = Inches(0.6)
table_y = Inches(1.4)
row_h = Inches(0.48)

# Header
add_rect(slide, table_x, table_y, Inches(8.5), Inches(0.5), C_PRIMARY)
for text, tx in [('評分項目', Inches(0.8)), ('權重', Inches(3.3)), ('說明', Inches(4.6))]:
    add_text_box(slide, tx, table_y + Inches(0.05), Inches(2.5), Inches(0.4),
                 text, font_size=16, color=C_WHITE, bold=True)

# Rows
for i, (name, weight, desc) in enumerate(categories):
    ry = table_y + Inches(0.5) + row_h * i
    bg = C_WHITE if i % 2 == 0 else C_LTGRAY
    add_rect(slide, table_x, ry, Inches(8.5), row_h, bg)
    add_text_box(slide, Inches(0.8), ry + Inches(0.07), Inches(2.3), Inches(0.35),
                 name, font_size=14, color=C_BLACK, bold=True)
    add_text_box(slide, Inches(3.3), ry + Inches(0.07), Inches(1.2), Inches(0.35),
                 weight, font_size=14, color=C_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.6), ry + Inches(0.07), Inches(4.3), Inches(0.35),
                 desc, font_size=14, color=C_GRAY)

# Right: Design insight
ix = Inches(9.5)
add_rect(slide, ix, Inches(1.4), Inches(3.2), Inches(5.7), C_WHITE)
add_rect(slide, ix, Inches(1.4), Inches(0.07), Inches(5.7), C_TEAL)

add_text_box(slide, ix + Inches(0.25), Inches(1.55), Inches(2.7), Inches(0.5),
             '🎯 設計理念', font_size=20, color=C_TEAL, bold=True)

insights = [
    '沒有「滿分策略」— 每個決策都有取捨',
    '能源＋韌性佔 45% — 強調供電核心目標',
    '外部成本 5% — 但抗議過高會觸發停工',
    '時間獎勵 10% — 鼓勵決策效率',
    '呼應 SDG 7、11、13'
]
add_multiline_box(slide, ix + Inches(0.25), Inches(2.2), Inches(2.7), Inches(4.5),
                  insights, font_size=14, color=C_BLACK, line_space=2.0)


# ════════════════════════════════════════════
# SLIDE 7: 台灣生質能現況
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
section_header(slide, '06', '現實連結 — 台灣生質能發展現況')

# Stats row
stat_block(slide, Inches(0.4),  Inches(1.4), '7.7%', '生質能占再生能源', C_PRIMARY)
stat_block(slide, Inches(3.3),  Inches(1.4), '850萬噸', '年廚餘產生量', C_ORANGE)
stat_block(slide, Inches(6.2),  Inches(1.4), '160+', '沼氣發電廠', C_TEAL)
stat_block(slide, Inches(9.5),  Inches(1.4), '2050', '淨零排放目標年', C_SEC)

# 3 Cases
cases = [
    ('🐄 畜牧沼氣發電', C_PRIMARY, [
        '全台約 160 座養豬場沼氣發電',
        '豬糞尿厭氧消化產生沼氣',
        '減少溫室氣體排放並發電'
    ]),
    ('🏙️ 都市廚餘發電', C_ORANGE, [
        '台北市生質能中心（外雙溪）',
        '每日處理 200 噸廚餘',
        '厭氧消化產沼氣發電'
    ]),
    ('🌾 農業廢棄物利用', C_TEAL, [
        '稻稈、蔗渣等氣化或發電',
        '花蓮光華園區案例',
        '結合在地農業的循環經濟'
    ])
]
for i, (title, color, items) in enumerate(cases):
    cx = Inches(0.6 + i * 4.2)
    add_rect(slide, cx, Inches(2.8), Inches(3.9), Inches(2.5), C_WHITE)
    add_rect(slide, cx, Inches(2.8), Inches(0.07), Inches(2.5), color)
    add_text_box(slide, cx + Inches(0.25), Inches(2.9), Inches(3.4), Inches(0.45),
                 title, font_size=18, color=color, bold=True)
    add_multiline_box(slide, cx + Inches(0.25), Inches(3.45), Inches(3.4), Inches(1.7),
                      items, font_size=14, color=C_BLACK, line_space=1.6)

# SDGs
add_rect(slide, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.6), C_WHITE)
add_rect(slide, Inches(0.6), Inches(5.6), Inches(0.07), Inches(1.6), C_PRIMARY)
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(11.4), Inches(0.5),
             '🌍 連結聯合國永續發展目標 SDGs', font_size=20, color=C_PRIMARY, bold=True)
add_multiline_box(slide, Inches(1.0), Inches(6.25), Inches(11.4), Inches(0.9), [
    'SDG 7 可負擔的潔淨能源 → 生質能是重要的再生能源來源',
    'SDG 11 永續城鄉 → 微電網提升社區能源韌性',
    'SDG 12 負責任消費與生產 → 廢棄物轉化為資源的循環經濟',
    'SDG 13 氣候行動 → 減少化石燃料依賴與溫室氣體排放'
], font_size=14, color=C_BLACK, line_space=1.4)


# ════════════════════════════════════════════
# SLIDE 8: 反思與延伸討論
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_ACCENT)
section_header(slide, '07', '反思與延伸討論')

questions = [
    ('Q1', '如果遊戲增加「太陽能板」和「小型風機」選項，你的策略會怎麼改變？生質能與其他再生能源如何互補？', C_PRIMARY),
    ('Q2', '遊戲中的「抗議值」對應現實中的什麼？如果你是能源開發商，會如何與社區溝通？', C_ORANGE),
    ('Q3', '台灣 2050 淨零排放目標中，生質能可以扮演什麼角色？它的限制是什麼？', C_TEAL),
    ('Q4', '你的學校有哪些廢棄物可以用來發電？請設計一個校園生質能系統的初步方案。', C_SEC),
]

for i, (qnum, qtext, color) in enumerate(questions):
    qy = Inches(1.4 + i * 1.45)
    add_rect(slide, Inches(0.6), qy, Inches(12.1), Inches(1.25), C_WHITE)
    add_rect(slide, Inches(0.6), qy, Inches(0.07), Inches(1.25), color)

    q_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), qy + Inches(0.3), Inches(0.6), Inches(0.6))
    q_shape.fill.solid()
    q_shape.fill.fore_color.rgb = color
    q_shape.line.fill.background()
    tf = q_shape.text_frame
    p = tf.paragraphs[0]
    p.text = qnum
    p.font.size = Pt(14)
    p.font.color.rgb = C_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.9), qy + Inches(0.25), Inches(10.5), Inches(0.8),
                 qtext, font_size=18, color=C_BLACK, line_space=1.4)


# ════════════════════════════════════════════
# SLIDE 9: Ending
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_DARK)
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.08), C_SEC)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.0),
             '🌱 能源轉型，從校園開始',
             font_size=44, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER, font_name=FONT_TITLE)

add_rect(slide, Inches(5.5), Inches(3.0), Inches(2.3), Inches(0.05), C_SEC)

add_text_box(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.8),
             '每一個設計決策，都在為永續未來投票',
             font_size=24, color=C_SEC, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.6), Inches(11.3), Inches(0.6),
             '請完成學習單，寫下你的反思與收穫！',
             font_size=20, color=C_LTGRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(3), Inches(5.5), Inches(7.3), Inches(1.1), C_PRIMARY)
add_text_box(slide, Inches(3.3), Inches(5.7), Inches(6.7), Inches(0.7),
             '📝 作業提醒：完成「課後反思學習單」並於下週繳交',
             font_size=18, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), SH - Inches(0.08), SW, Inches(0.08), C_SEC)


# ── Save ──
prs.save(OUTPUT)
print(f'PPTX created: {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')
